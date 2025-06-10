"""
数据库模型和管理
Database models and management for Ruishi Control Platform
"""

import os
import sqlite3
import hashlib
import json
from datetime import datetime
from typing import List, Dict, Any, Optional
import logging

logger = logging.getLogger(__name__)

class DatabaseManager:
    """数据库管理器"""
    
    def __init__(self, db_path: str = None):
        if db_path is None:
            db_path = os.path.join(os.path.dirname(__file__), '..', 'data', 'ruishi_platform.db')
        
        self.db_path = db_path
        self.ensure_data_directory()
        self.init_database()
    
    def ensure_data_directory(self):
        """确保数据目录存在"""
        data_dir = os.path.dirname(self.db_path)
        if not os.path.exists(data_dir):
            os.makedirs(data_dir)
    
    def get_connection(self):
        """获取数据库连接"""
        conn = sqlite3.connect(self.db_path)
        conn.row_factory = sqlite3.Row  # 使结果可以通过列名访问
        return conn
    
    def init_database(self):
        """初始化数据库表"""
        conn = self.get_connection()
        try:
            # 创建用户表
            conn.execute('''
                CREATE TABLE IF NOT EXISTS users (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    username TEXT UNIQUE NOT NULL,
                    password_hash TEXT NOT NULL,
                    email TEXT,
                    role TEXT DEFAULT 'user',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    last_login TIMESTAMP,
                    is_active BOOLEAN DEFAULT 1
                )
            ''')
            
            # 创建文档表
            conn.execute('''
                CREATE TABLE IF NOT EXISTS documents (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    filename TEXT NOT NULL,
                    original_filename TEXT NOT NULL,
                    file_path TEXT NOT NULL,
                    file_size INTEGER,
                    file_type TEXT,
                    mime_type TEXT,
                    category TEXT DEFAULT 'general',
                    title TEXT,
                    description TEXT,
                    content_text TEXT,
                    content_summary TEXT,
                    keywords TEXT,
                    uploaded_by INTEGER,
                    upload_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    last_modified TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    is_active BOOLEAN DEFAULT 1,
                    download_count INTEGER DEFAULT 0,
                    FOREIGN KEY (uploaded_by) REFERENCES users (id)
                )
            ''')
            
            # 创建文档标签表
            conn.execute('''
                CREATE TABLE IF NOT EXISTS document_tags (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    document_id INTEGER,
                    tag TEXT,
                    FOREIGN KEY (document_id) REFERENCES documents (id),
                    UNIQUE(document_id, tag)
                )
            ''')
            
            # 创建用户会话表
            conn.execute('''
                CREATE TABLE IF NOT EXISTS user_sessions (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    user_id INTEGER,
                    session_token TEXT UNIQUE,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    expires_at TIMESTAMP,
                    is_active BOOLEAN DEFAULT 1,
                    FOREIGN KEY (user_id) REFERENCES users (id)
                )
            ''')
            
            # 创建AI问答历史表
            conn.execute('''
                CREATE TABLE IF NOT EXISTS ai_conversations (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    user_id INTEGER,
                    session_id TEXT,
                    question TEXT NOT NULL,
                    answer TEXT,
                    ai_provider TEXT,
                    ai_model TEXT,
                    related_documents TEXT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    response_time REAL,
                    rating INTEGER,
                    FOREIGN KEY (user_id) REFERENCES users (id)
                )
            ''')
            
            # 创建系统配置表
            conn.execute('''
                CREATE TABLE IF NOT EXISTS system_config (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    config_key TEXT UNIQUE,
                    config_value TEXT,
                    description TEXT,
                    updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            
            conn.commit()
            
            # 创建默认管理员账户
            self.create_default_admin()
            
            logger.info("数据库初始化完成")
            
        except Exception as e:
            logger.error(f"数据库初始化失败: {e}")
            conn.rollback()
        finally:
            conn.close()
    
    def create_default_admin(self):
        """创建默认管理员账户"""
        conn = self.get_connection()
        try:
            # 检查是否已存在管理员
            cursor = conn.execute("SELECT id FROM users WHERE username = 'admin'")
            if cursor.fetchone():
                return
            
            # 创建默认管理员账户
            password_hash = self.hash_password('admin123')
            conn.execute('''
                INSERT INTO users (username, password_hash, email, role)
                VALUES (?, ?, ?, ?)
            ''', ('admin', password_hash, 'admin@jytek.com', 'admin'))
            
            conn.commit()
            logger.info("默认管理员账户创建成功 - 用户名: admin, 密码: admin123")
            
        except Exception as e:
            logger.error(f"创建默认管理员失败: {e}")
            conn.rollback()
        finally:
            conn.close()
    
    @staticmethod
    def hash_password(password: str) -> str:
        """密码哈希"""
        return hashlib.sha256(password.encode()).hexdigest()
    
    def verify_password(self, password: str, password_hash: str) -> bool:
        """验证密码"""
        return self.hash_password(password) == password_hash


class UserManager:
    """用户管理器"""
    
    def __init__(self, db_manager: DatabaseManager):
        self.db = db_manager
    
    def authenticate(self, username: str, password: str) -> Optional[Dict]:
        """用户认证"""
        conn = self.db.get_connection()
        try:
            cursor = conn.execute('''
                SELECT id, username, password_hash, email, role, is_active
                FROM users WHERE username = ? AND is_active = 1
            ''', (username,))
            
            user = cursor.fetchone()
            if user and self.db.verify_password(password, user['password_hash']):
                # 更新最后登录时间
                conn.execute('''
                    UPDATE users SET last_login = CURRENT_TIMESTAMP WHERE id = ?
                ''', (user['id'],))
                conn.commit()
                
                return {
                    'id': user['id'],
                    'username': user['username'],
                    'email': user['email'],
                    'role': user['role']
                }
            return None
            
        except Exception as e:
            logger.error(f"用户认证失败: {e}")
            return None
        finally:
            conn.close()
    
    def create_session(self, user_id: int) -> str:
        """创建用户会话"""
        import uuid
        from datetime import timedelta
        
        session_token = str(uuid.uuid4())
        expires_at = datetime.now() + timedelta(hours=24)  # 24小时过期
        
        conn = self.db.get_connection()
        try:
            conn.execute('''
                INSERT INTO user_sessions (user_id, session_token, expires_at)
                VALUES (?, ?, ?)
            ''', (user_id, session_token, expires_at))
            conn.commit()
            return session_token
            
        except Exception as e:
            logger.error(f"创建会话失败: {e}")
            return None
        finally:
            conn.close()
    
    def verify_session(self, session_token: str) -> Optional[Dict]:
        """验证会话"""
        conn = self.db.get_connection()
        try:
            cursor = conn.execute('''
                SELECT s.user_id, u.username, u.email, u.role
                FROM user_sessions s
                JOIN users u ON s.user_id = u.id
                WHERE s.session_token = ? AND s.is_active = 1 
                AND s.expires_at > CURRENT_TIMESTAMP
            ''', (session_token,))
            
            result = cursor.fetchone()
            if result:
                return {
                    'id': result['user_id'],
                    'username': result['username'],
                    'email': result['email'],
                    'role': result['role']
                }
            return None
            
        except Exception as e:
            logger.error(f"会话验证失败: {e}")
            return None
        finally:
            conn.close()
    
    def logout(self, session_token: str):
        """用户登出"""
        conn = self.db.get_connection()
        try:
            conn.execute('''
                UPDATE user_sessions SET is_active = 0 WHERE session_token = ?
            ''', (session_token,))
            conn.commit()
            
        except Exception as e:
            logger.error(f"登出失败: {e}")
        finally:
            conn.close()


class DocumentManager:
    """文档管理器"""
    
    def __init__(self, db_manager: DatabaseManager):
        self.db = db_manager
        self.upload_dir = os.path.join(os.path.dirname(__file__), '..', 'data', 'uploads')
        self.ensure_upload_directory()
    
    def ensure_upload_directory(self):
        """确保上传目录存在"""
        if not os.path.exists(self.upload_dir):
            os.makedirs(self.upload_dir)
    
    def save_document(self, file_data: bytes, filename: str, user_id: int, 
                     category: str = 'general', title: str = None, 
                     description: str = None) -> Optional[int]:
        """保存文档"""
        import uuid
        import mimetypes
        
        try:
            # 生成唯一文件名
            file_ext = os.path.splitext(filename)[1]
            unique_filename = f"{uuid.uuid4()}{file_ext}"
            file_path = os.path.join(self.upload_dir, unique_filename)
            
            # 保存文件
            with open(file_path, 'wb') as f:
                f.write(file_data)
            
            # 获取文件信息
            file_size = len(file_data)
            mime_type, _ = mimetypes.guess_type(filename)
            file_type = self.get_file_type(filename)
            
            # 提取文本内容
            content_text = self.extract_text_content(file_path, file_type)
            content_summary = self.generate_summary(content_text)
            
            # 保存到数据库
            conn = self.db.get_connection()
            try:
                cursor = conn.execute('''
                    INSERT INTO documents (
                        filename, original_filename, file_path, file_size, 
                        file_type, mime_type, category, title, description,
                        content_text, content_summary, uploaded_by
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                ''', (
                    unique_filename, filename, file_path, file_size,
                    file_type, mime_type, category, title or filename, description,
                    content_text, content_summary, user_id
                ))
                
                document_id = cursor.lastrowid
                conn.commit()
                
                logger.info(f"文档保存成功: {filename} -> {unique_filename}")
                return document_id
                
            except Exception as e:
                logger.error(f"数据库保存失败: {e}")
                # 删除已保存的文件
                if os.path.exists(file_path):
                    os.remove(file_path)
                return None
            finally:
                conn.close()
                
        except Exception as e:
            logger.error(f"文档保存失败: {e}")
            return None
    
    def get_file_type(self, filename: str) -> str:
        """获取文件类型"""
        ext = os.path.splitext(filename)[1].lower()
        type_mapping = {
            '.pdf': 'pdf',
            '.doc': 'word',
            '.docx': 'word',
            '.txt': 'text',
            '.md': 'markdown',
            '.ppt': 'powerpoint',
            '.pptx': 'powerpoint',
            '.xls': 'excel',
            '.xlsx': 'excel',
            '.jpg': 'image',
            '.jpeg': 'image',
            '.png': 'image',
            '.gif': 'image'
        }
        return type_mapping.get(ext, 'unknown')
    
    def extract_text_content(self, file_path: str, file_type: str) -> str:
        """提取文本内容"""
        try:
            if file_type == 'text':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            elif file_type == 'pdf':
                return self.extract_pdf_text(file_path)
            elif file_type == 'word':
                return self.extract_word_text(file_path)
            elif file_type == 'powerpoint':
                return self.extract_ppt_text(file_path)
            else:
                return ""
        except Exception as e:
            logger.error(f"文本提取失败: {e}")
            return ""
    
    def extract_pdf_text(self, file_path: str) -> str:
        """提取PDF文本"""
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except ImportError:
            logger.warning("PyPDF2未安装，无法提取PDF文本")
            return ""
        except Exception as e:
            logger.error(f"PDF文本提取失败: {e}")
            return ""
    
    def extract_word_text(self, file_path: str) -> str:
        """提取Word文档文本"""
        try:
            import docx
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except ImportError:
            logger.warning("python-docx未安装，无法提取Word文本")
            return ""
        except Exception as e:
            logger.error(f"Word文本提取失败: {e}")
            return ""
    
    def extract_ppt_text(self, file_path: str) -> str:
        """提取PowerPoint文本"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            logger.warning("python-pptx未安装，无法提取PowerPoint文本")
            return ""
        except Exception as e:
            logger.error(f"PowerPoint文本提取失败: {e}")
            return ""
    
    def generate_summary(self, content: str, max_length: int = 200) -> str:
        """生成内容摘要"""
        if not content:
            return ""
        
        # 简单的摘要生成：取前200个字符
        summary = content.strip()[:max_length]
        if len(content) > max_length:
            summary += "..."
        
        return summary
    
    def search_documents(self, query: str, category: str = None, limit: int = 10) -> List[Dict]:
        """搜索文档"""
        conn = self.db.get_connection()
        try:
            sql = '''
                SELECT id, filename, original_filename, file_type, category,
                       title, description, content_summary, upload_time, file_size
                FROM documents 
                WHERE is_active = 1
            '''
            params = []
            
            if query:
                sql += ' AND (title LIKE ? OR description LIKE ? OR content_text LIKE ?)'
                query_param = f'%{query}%'
                params.extend([query_param, query_param, query_param])
            
            if category:
                sql += ' AND category = ?'
                params.append(category)
            
            sql += ' ORDER BY upload_time DESC LIMIT ?'
            params.append(limit)
            
            cursor = conn.execute(sql, params)
            results = []
            
            for row in cursor.fetchall():
                results.append({
                    'id': row['id'],
                    'filename': row['filename'],
                    'original_filename': row['original_filename'],
                    'file_type': row['file_type'],
                    'category': row['category'],
                    'title': row['title'],
                    'description': row['description'],
                    'content_summary': row['content_summary'],
                    'upload_time': row['upload_time'],
                    'file_size': row['file_size']
                })
            
            return results
            
        except Exception as e:
            logger.error(f"文档搜索失败: {e}")
            return []
        finally:
            conn.close()
    
    def get_document_content(self, document_id: int) -> Optional[str]:
        """获取文档内容"""
        conn = self.db.get_connection()
        try:
            cursor = conn.execute('''
                SELECT content_text FROM documents WHERE id = ? AND is_active = 1
            ''', (document_id,))
            
            result = cursor.fetchone()
            return result['content_text'] if result else None
            
        except Exception as e:
            logger.error(f"获取文档内容失败: {e}")
            return None
        finally:
            conn.close()
    
    def get_all_documents(self, page: int = 1, per_page: int = 20) -> Dict:
        """获取所有文档（分页）"""
        conn = self.db.get_connection()
        try:
            offset = (page - 1) * per_page
            
            # 获取总数
            cursor = conn.execute('SELECT COUNT(*) as total FROM documents WHERE is_active = 1')
            total = cursor.fetchone()['total']
            
            # 获取文档列表
            cursor = conn.execute('''
                SELECT d.*, u.username as uploaded_by_name
                FROM documents d
                LEFT JOIN users u ON d.uploaded_by = u.id
                WHERE d.is_active = 1
                ORDER BY d.upload_time DESC
                LIMIT ? OFFSET ?
            ''', (per_page, offset))
            
            documents = []
            for row in cursor.fetchall():
                documents.append({
                    'id': row['id'],
                    'filename': row['filename'],
                    'original_filename': row['original_filename'],
                    'file_type': row['file_type'],
                    'category': row['category'],
                    'title': row['title'],
                    'description': row['description'],
                    'content_summary': row['content_summary'],
                    'upload_time': row['upload_time'],
                    'file_size': row['file_size'],
                    'uploaded_by_name': row['uploaded_by_name'],
                    'download_count': row['download_count']
                })
            
            return {
                'documents': documents,
                'total': total,
                'page': page,
                'per_page': per_page,
                'pages': (total + per_page - 1) // per_page
            }
            
        except Exception as e:
            logger.error(f"获取文档列表失败: {e}")
            return {'documents': [], 'total': 0, 'page': 1, 'per_page': per_page, 'pages': 0}
        finally:
            conn.close()
    
    def delete_document(self, document_id: int) -> bool:
        """删除文档"""
        conn = self.db.get_connection()
        try:
            # 获取文件路径
            cursor = conn.execute('SELECT file_path FROM documents WHERE id = ?', (document_id,))
            result = cursor.fetchone()
            
            if result:
                # 软删除（标记为不活跃）
                conn.execute('UPDATE documents SET is_active = 0 WHERE id = ?', (document_id,))
                conn.commit()
                
                # 可选：删除物理文件
                # if os.path.exists(result['file_path']):
                #     os.remove(result['file_path'])
                
                return True
            return False
            
        except Exception as e:
            logger.error(f"删除文档失败: {e}")
            return False
        finally:
            conn.close()


# 全局数据库管理器实例
db_manager = DatabaseManager()
user_manager = UserManager(db_manager)
document_manager = DocumentManager(db_manager)
